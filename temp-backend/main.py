import base64
import json
import logging
from typing import List, Optional, Union
from fastapi import FastAPI, UploadFile, File, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from fastapi.security import OAuth2PasswordBearer
from pydantic import BaseModel
import httpx
from pypdf import PdfReader
from sqlmodel import SQLModel, Field, Session, create_engine, select
import bcrypt
import os
from jose import JWTError, jwt
from datetime import datetime, timedelta, timezone

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="OmniAI Backend")

# Setup CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Adjust in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

DATABASE_URL = os.environ.get("DATABASE_URL")

if DATABASE_URL and DATABASE_URL.startswith("postgres://"):
    DATABASE_URL = DATABASE_URL.replace("postgres://", "postgresql://", 1)

sqlite_url = "sqlite:///./omniai.db"
engine = create_engine(DATABASE_URL or sqlite_url, connect_args={"check_same_thread": False} if not DATABASE_URL else {})

SECRET_KEY = os.environ.get("SECRET_KEY", "super_secret_omniai_key_for_testing")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 24 * 7

class User(SQLModel, table=True):
    __table_args__ = {'extend_existing': True}
    id: Optional[int] = Field(default=None, primary_key=True)
    username: str = Field(unique=True, index=True)
    hashed_password: str

class UserCreate(BaseModel):
    username: str
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str

def get_password_hash(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(plain_password: str, hashed_password: str) -> bool:
    return bcrypt.checkpw(plain_password.encode('utf-8'), hashed_password.encode('utf-8'))

def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

SQLModel.metadata.create_all(engine)

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="auth/login")

def get_current_user(token: str = Depends(oauth2_scheme)):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise HTTPException(status_code=401, detail="Invalid token")
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid token")
    
    with Session(engine) as session:
        user = session.exec(select(User).where(User.username == username)).first()
        if user is None:
            raise HTTPException(status_code=401, detail="User not found")
        return user

class ChatMessage(BaseModel):
    role: str
    content: Union[str, List[dict]]

class ChatRequest(BaseModel):
    base_url: str
    api_key: str
    model: str
    messages: List[ChatMessage]

class WebSearchRequest(BaseModel):
    query: str
    max_results: Optional[int] = 5

class ConnectRequest(BaseModel):
    provider_name: str = ""
    base_url: str
    api_key: str
    model: str = ""

@app.post("/upload-image")
async def upload_image(file: UploadFile = File(...)):
    if not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="Uploaded file is not an image.")
    
    try:
        contents = await file.read()
        encoded = base64.b64encode(contents).decode("utf-8")
        return {
            "filename": file.filename,
            "mime_type": file.content_type,
            "base64": encoded
        }
    except Exception as e:
        logger.error(f"Error uploading image: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process image: {str(e)}")

def extract_text_from_bytes(filename: str, content_type: str, contents: bytes) -> str:
    filename_lower = filename.lower()
    
    if filename_lower.endswith(".pdf") or content_type == "application/pdf":
        from io import BytesIO
        pdf = PdfReader(BytesIO(contents))
        extracted_text = []
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                extracted_text.append(page_text)
        return "\n".join(extracted_text)
        
    elif filename_lower.endswith(".docx") or content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import docx
        from io import BytesIO
        doc = docx.Document(BytesIO(contents))
        text = "\n".join([para.text for para in doc.paragraphs])
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_texts.append(" | ".join(row_text))
        if table_texts:
            text += "\n\nTables:\n" + "\n".join(table_texts)
        return text
        
    elif filename_lower.endswith(".xlsx") or content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        import openpyxl
        from io import BytesIO
        wb = openpyxl.load_workbook(BytesIO(contents), data_only=True)
        sheets_text = []
        for name in wb.sheetnames:
            sheet = wb[name]
            sheets_text.append(f"--- Sheet: {name} ---")
            for row in sheet.iter_rows(values_only=True):
                if any(row):
                    row_str = "\t".join([str(val) if val is not None else "" for val in row])
                    sheets_text.append(row_str)
        return "\n".join(sheets_text)
        
    elif filename_lower.endswith(".pptx") or content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        from pptx import Presentation
        from io import BytesIO
        prs = Presentation(BytesIO(contents))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slides_text.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slides_text.append(shape.text)
        return "\n".join(slides_text)
        
    elif filename_lower.endswith(".zip") or content_type in ["application/zip", "application/x-zip-compressed", "application/octet-stream"]:
        import zipfile
        from io import BytesIO
        zip_text = []
        with zipfile.ZipFile(BytesIO(contents)) as z:
            for name in z.namelist():
                if name.endswith("/") or name.startswith("__MACOSX/") or ".DS_Store" in name:
                    continue
                name_lower = name.lower()
                sub_mime = ""
                if name_lower.endswith(".pdf"):
                    sub_mime = "application/pdf"
                elif name_lower.endswith(".docx"):
                    sub_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                elif name_lower.endswith(".xlsx"):
                    sub_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                elif name_lower.endswith(".pptx"):
                    sub_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                
                try:
                    file_bytes = z.read(name)
                    if name_lower.endswith((".txt", ".md", ".csv", ".json", ".xml", ".html", ".pdf", ".docx", ".xlsx", ".pptx")):
                        extracted = extract_text_from_bytes(name, sub_mime, file_bytes)
                        zip_text.append(f"=== File: {name} ===\n{extracted}\n")
                except Exception as fe:
                    zip_text.append(f"=== File: {name} (Error parsing: {str(fe)}) ===\n")
        return "\n".join(zip_text)
        
    else:
        try:
            return contents.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return contents.decode("latin-1")
            except Exception as e:
                return f"[Error decoding text content: {str(e)}]"

@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    try:
        content_type = file.content_type or ""
        filename = file.filename or ""
        
        contents = await file.read()
        text = extract_text_from_bytes(filename, content_type, contents)
                
        return {
            "filename": filename,
            "text": text
        }
    except Exception as e:
        logger.error(f"Error uploading file: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process file: {str(e)}")

async def stream_chat_response(request: ChatRequest):
    headers = {
        "Authorization": f"Bearer {request.api_key}",
        "Content-Type": "application/json",
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
    }
    
    # Standardize messages to JSON serializable objects
    serialized_messages = []
    for msg in request.messages:
        serialized_messages.append({
            "role": msg.role,
            "content": msg.content
        })
        
    payload = {
        "model": request.model,
        "messages": serialized_messages,
        "stream": True
    }
    
    clean_url = request.base_url.rstrip("/")
    if not clean_url.endswith("/v1") and "openrouter.ai" in clean_url:
        target_url = f"{clean_url}/api/v1/chat/completions"
    elif not clean_url.endswith("/chat/completions"):
        target_url = f"{clean_url}/chat/completions"
    else:
        target_url = clean_url

    logger.info(f"Forwarding request to {target_url} with model {request.model}")
    
    async with httpx.AsyncClient() as client:
        try:
            async with client.stream(
                "POST",
                target_url,
                headers=headers,
                json=payload,
                timeout=60.0
            ) as response:
                if response.status_code != 200:
                    error_detail = await response.aread()
                    logger.error(f"Provider API returned error status {response.status_code}: {error_detail}")
                    yield f"Error contacting AI: Provider returned status {response.status_code}"
                    return
                
                buffer = ""
                async for line in response.aiter_lines():
                    if line.startswith("data: "):
                        data_str = line[6:].strip()
                        if data_str == "[DONE]":
                            continue
                        try:
                            data = json.loads(data_str)
                            delta = data.get("choices", [{}])[0].get("delta", {})
                            if "content" in delta and delta["content"]:
                                yield delta["content"]
                        except Exception:
                            pass
                    else:
                        buffer += line
                        
                # Fallback if API completely ignored stream=True and returned strict JSON block
                if buffer.strip() and buffer.strip().startswith("{") and buffer.strip().endswith("}"):
                    try:
                        data = json.loads(buffer)
                        content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
                        if content:
                            yield content
                    except Exception:
                        pass
                        
        except Exception as e:
            logger.error(f"Exception during AI completions streaming: {e.__class__.__name__} - {str(e)}")
            yield f"\n[Backend Error: Provider disconnected unexpectedly: {e.__class__.__name__}]"

@app.post("/web-search")
async def web_search(request: WebSearchRequest, current_user: User = Depends(get_current_user)):
    from duckduckgo_search import DDGS
    try:
        with DDGS() as ddgs:
            raw = list(ddgs.text(request.query, max_results=request.max_results))
            results = [
                {
                    "title": r.get("title", ""),
                    "url": r.get("href", r.get("link", "")),
                    "snippet": r.get("body", r.get("snippet", "")),
                }
                for r in raw
            ]
        return {"results": results}
    except Exception as e:
        logger.error(f"Web search error: {e}")
        return {"results": [], "error": str(e)}

@app.post("/auth/register")
def register(user: UserCreate):
    with Session(engine) as session:
        statement = select(User).where(User.username == user.username)
        if session.exec(statement).first():
            raise HTTPException(status_code=400, detail="Username already registered")
        db_user = User(username=user.username, hashed_password=get_password_hash(user.password))
        session.add(db_user)
        session.commit()
    return {"message": "User registered successfully"}

@app.post("/auth/login", response_model=Token)
def login(user: UserCreate):
    with Session(engine) as session:
        statement = select(User).where(User.username == user.username)
        db_user = session.exec(statement).first()
        if not db_user or not verify_password(user.password, db_user.hashed_password):
            raise HTTPException(status_code=401, detail="Incorrect username or password")
        
        access_token = create_access_token(data={"sub": db_user.username})
        return {"access_token": access_token, "token_type": "bearer"}

@app.post("/connect")
async def connect_provider(req: ConnectRequest, current_user: User = Depends(get_current_user)):
    clean_url = req.base_url.rstrip("/")
    if "openrouter" in clean_url and not clean_url.endswith("/models"):
        models_url = f"{clean_url}/models"
    else:
        models_url = f"{clean_url}/models"
        
    headers = {
        "Authorization": f"Bearer {req.api_key}",
        "Content-Type": "application/json",
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
    }
    
    async with httpx.AsyncClient() as client:
        try:
            res = await client.get(models_url, headers=headers, timeout=10.0)
            res.raise_for_status()
            
            try:
                data = res.json()
                if isinstance(data, list):
                    raw_models = data
                elif isinstance(data, dict):
                    raw_models = data.get("data", data.get("models", []))
                    if not isinstance(raw_models, list):
                        raw_models = [raw_models] if raw_models else []
                else:
                    raw_models = []
            except Exception:
                raw_models = []
                
            models_list = []
            for m in raw_models:
                if not isinstance(m, dict): continue
                models_list.append({
                    "id": m.get("id", ""),
                    "name": m.get("name", m.get("id", "")),
                    "description": m.get("description", ""),
                    "context_length": m.get("context_length", 0),
                    "free": m.get("pricing", {}).get("prompt", "1") == "0" if isinstance(m.get("pricing"), dict) else False
                })
            
            return {"models": models_list}
        except Exception as e:
            logger.error(f"Failed to parse models from {models_url}: {e}")
            # Do not crash the app. Allow them to proceed with a default fallback model.
            generic_models = [{"id": "fallback-model", "name": "Any Model (Manual Entry Supported)", "description": "Type your model manually", "context_length": 8192, "free": True}]
            return {"message": "Connected with fallback", "models": generic_models}

@app.post("/chat")
async def chat(request: ChatRequest, current_user: User = Depends(get_current_user)):
    return StreamingResponse(stream_chat_response(request), media_type="text/event-stream")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)
