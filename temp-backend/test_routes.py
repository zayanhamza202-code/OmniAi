import os
import requests
import io
import time
import zipfile
import tempfile

API_URL = "http://127.0.0.1:8000"

# ─── Phase 1 Tests ───────────────────────────────────────────────────

def test_upload_image():
    print("Testing /upload-image endpoint...")
    dummy_img = (
        b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06'
        b'\x00\x00\x00\x1f\x15c4\x00\x00\x00\nIDATx\x9cc`\x00\x00\x00\x02\x00\x01H\xaf'
        b'\xa4q\x00\x00\x00\x00IEND\xaeB`\x82'
    )
    files = {"file": ("test.png", dummy_img, "image/png")}
    response = requests.post(f"{API_URL}/upload-image", files=files)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert data["mime_type"] == "image/png"
    assert len(data["base64"]) > 0
    print("SUCCESS: /upload-image works!")

def test_upload_text():
    print("\nTesting /upload (text) endpoint...")
    dummy_text = b"Hello, this is a test text file for OmniAI."
    files = {"file": ("test.txt", dummy_text, "text/plain")}
    response = requests.post(f"{API_URL}/upload", files=files)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert data["text"] == "Hello, this is a test text file for OmniAI."
    print("SUCCESS: /upload (text) works!")

def test_chat_proxy():
    print("\nTesting /chat endpoint (streaming error handling)...")
    payload = {
        "base_url": "https://openrouter.ai/api/v1",
        "api_key": "dummy_key",
        "model": "google/gemini-flash-1.5",
        "messages": [{"role": "user", "content": "Hello"}]
    }
    response = requests.post(f"{API_URL}/chat", json=payload, stream=True)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    output = b""
    for chunk in response.iter_content(chunk_size=1024):
        if chunk:
            output += chunk
    decoded = output.decode("utf-8")
    print(f"SUCCESS: /chat streamed: {decoded[:80]}...")

# ─── Phase 2 Tests ───────────────────────────────────────────────────

def test_upload_csv():
    print("\nTesting /upload (CSV)...")
    csv_data = b"Name,Age,City\nAlice,30,NYC\nBob,25,London\n"
    files = {"file": ("data.csv", csv_data, "text/csv")}
    response = requests.post(f"{API_URL}/upload", files=files)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert "Alice" in data["text"]
    assert "Bob" in data["text"]
    assert data["filename"] == "data.csv"
    print(f"SUCCESS: CSV extracted ({len(data['text'])} chars)")

def test_upload_markdown():
    print("\nTesting /upload (Markdown)...")
    md_data = b"# Hello World\n\nThis is **bold** and *italic* text.\n\n- Item 1\n- Item 2\n"
    files = {"file": ("readme.md", md_data, "text/markdown")}
    response = requests.post(f"{API_URL}/upload", files=files)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert "Hello World" in data["text"]
    assert "bold" in data["text"]
    assert data["filename"] == "readme.md"
    print(f"SUCCESS: Markdown extracted ({len(data['text'])} chars)")

def test_upload_zip():
    print("\nTesting /upload (ZIP)...")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("notes/hello.txt", "Hello from inside the zip!")
        zf.writestr("data.csv", "x,y\n1,2\n3,4\n")
    buf.seek(0)
    files = {"file": ("archive.zip", buf.read(), "application/zip")}
    response = requests.post(f"{API_URL}/upload", files=files)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert "Hello from inside the zip!" in data["text"]
    assert "notes/hello.txt" in data["text"]
    assert data["filename"] == "archive.zip"
    print(f"SUCCESS: ZIP extracted ({len(data['text'])} chars)")

def test_upload_docx():
    print("\nTesting /upload (DOCX)...")
    import docx
    doc = docx.Document()
    doc.add_paragraph("OmniAI DOCX Test Paragraph")
    doc.add_paragraph("Second paragraph with important data.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A1"
    table.cell(0, 1).text = "B1"
    table.cell(1, 0).text = "A2"
    table.cell(1, 1).text = "B2"
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    files = {"file": ("test.docx", buf.read(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
    response = requests.post(f"{API_URL}/upload", files=files)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert "OmniAI DOCX Test Paragraph" in data["text"]
    assert "Second paragraph" in data["text"]
    assert "A1" in data["text"]
    assert data["filename"] == "test.docx"
    print(f"SUCCESS: DOCX extracted ({len(data['text'])} chars)")

def test_upload_xlsx():
    print("\nTesting /upload (XLSX)...")
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "TestSheet"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    ws.append(["Bob", 87])
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    files = {"file": ("scores.xlsx", buf.read(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}
    response = requests.post(f"{API_URL}/upload", files=files)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert "TestSheet" in data["text"]
    assert "Alice" in data["text"]
    assert "95" in data["text"]
    assert data["filename"] == "scores.xlsx"
    print(f"SUCCESS: XLSX extracted ({len(data['text'])} chars)")

def test_upload_pptx():
    print("\nTesting /upload (PPTX)...")
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = "OmniAI Slide Title"
    slide.placeholders[1].text = "Bullet point content here"
    slide2 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide2.shapes.title.text = "Second Slide"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    files = {"file": ("deck.pptx", buf.read(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    response = requests.post(f"{API_URL}/upload", files=files)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert "OmniAI Slide Title" in data["text"]
    assert "Bullet point" in data["text"]
    assert "Second Slide" in data["text"]
    assert data["filename"] == "deck.pptx"
    print(f"SUCCESS: PPTX extracted ({len(data['text'])} chars)")

def test_web_search():
    print("\nTesting /web-search...")
    payload = {"query": "Python programming language", "max_results": 3}
    response = requests.post(f"{API_URL}/web-search", json=payload)
    assert response.status_code == 200, f"FAILED: status {response.status_code}"
    data = response.json()
    assert "results" in data
    assert len(data["results"]) > 0
    for r in data["results"]:
        assert "title" in r
        assert "url" in r
        assert "snippet" in r
    print(f"SUCCESS: Web search returned {len(data['results'])} results")
    for r in data["results"]:
        print(f"  - {r['title'][:60]}")

# ─── Runner ───────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("=" * 60)
    print("  OmniAI Backend Route Tests")
    print("=" * 60)

    tests = [
        ("Phase 1", test_upload_image),
        ("Phase 1", test_upload_text),
        ("Phase 1", test_chat_proxy),
        ("Phase 2", test_upload_csv),
        ("Phase 2", test_upload_markdown),
        ("Phase 2", test_upload_zip),
        ("Phase 2", test_upload_docx),
        ("Phase 2", test_upload_xlsx),
        ("Phase 2", test_upload_pptx),
        ("Phase 4", test_web_search),
    ]

    passed = 0
    failed = 0
    for phase, test_fn in tests:
        try:
            test_fn()
            passed += 1
        except Exception as e:
            print(f"FAILED: {test_fn.__name__} - {e}")
            failed += 1

    print("\n" + "=" * 60)
    print(f"  Results: {passed} passed, {failed} failed out of {len(tests)}")
    print("=" * 60)
